from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from patent_agent.core.models import EvidenceScope


@dataclass(frozen=True)
class PdfLogicalBlock:
    heading: str
    text: str
    page: int
    paragraph_index: int
    block_type: str = "paragraph"
    scope: EvidenceScope = EvidenceScope.INVENTION_SOURCE


def read_text(path: Path) -> list[tuple[str, str]]:
    return [("全文", path.read_text(encoding="utf-8"))]


def read_docx(path: Path) -> list[tuple[str, str]]:
    from docx import Document
    document = Document(path)
    blocks = []
    heading = "正文"
    buffer: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            if buffer: blocks.append((heading, "\n".join(buffer))); buffer = []
            heading = text
        else:
            buffer.append(text)
    if buffer: blocks.append((heading, "\n".join(buffer)))
    return blocks


def read_pdf(path: Path) -> list[tuple[str, str]]:
    return [(f"第{item.page}页 · {item.heading}", item.text) for item in read_pdf_blocks(path)]


def read_pdf_blocks(path: Path) -> list[PdfLogicalBlock]:
    from pypdf import PdfReader
    blocks: list[PdfLogicalBlock] = []
    paragraph_index = 0
    for page_number, page in enumerate(PdfReader(str(path)).pages, 1):
        for block in split_pdf_page_text(page.extract_text() or "", page_number):
            paragraph_index += 1
            blocks.append(PdfLogicalBlock(block.heading, block.text, page_number, paragraph_index, block.block_type, block.scope))
    return blocks


_ROMAN_HEADING = re.compile(r"^(?:[IVX]+)\.\s+(.+)$", re.I)
_LETTER_HEADING = re.compile(r"^[A-Z]\.\s+(.+)$")
_NUMBER_HEADING = re.compile(r"^\d+\)\s+(.+)$")
_REFERENCE = re.compile(r"^\[\d+\]\s+")
_FIGURE = re.compile(r"^Fig\.\s*\d+\.?\s*", re.I)
_TABLE = re.compile(r"^TABLE\s+[IVX]+", re.I)
_EQUATION = re.compile(r"(?:[=λΣ∫]|\\sum|\\int).*(?:\(\s*[\d-]+\s*\))\s*$")


def split_pdf_page_text(text: str, page_number: int) -> list[PdfLogicalBlock]:
    """Split selectable academic PDF text into semantic blocks, not fixed-size slices."""
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    output: list[PdfLogicalBlock] = []
    section = "Title" if page_number == 1 else f"Page {page_number}"
    scope = EvidenceScope.INVENTION_SOURCE
    block_type = "paragraph"
    buffer: list[str] = []
    local_index = 0

    def flush(kind: str | None = None) -> None:
        nonlocal buffer, local_index
        content = " ".join(buffer).strip()
        if content:
            local_index += 1
            output.append(PdfLogicalBlock(section, content, page_number, local_index, kind or block_type, scope))
        buffer = []

    for line in lines:
        upper = line.upper()
        if upper in {"VI. REFERENCES", "REFERENCES"} or re.match(r"^[IVX]+\.\s+REFERENCES$", upper):
            flush()
            section, scope, block_type = "References", EvidenceScope.REFERENCE, "reference"
            continue
        if scope == EvidenceScope.REFERENCE and _REFERENCE.match(line):
            flush("reference")
            buffer = [line]
            continue
        if scope == EvidenceScope.REFERENCE:
            buffer.append(line)
            continue
        if line.startswith("Abstract—") or line.startswith("Abstract-"):
            flush()
            section, block_type = "Abstract", "abstract"
            buffer = [re.sub(r"^Abstract[—-]", "", line).strip()]
            continue
        if line.startswith("Keywords—") or line.startswith("Keywords-"):
            flush()
            section, block_type = "Keywords", "keywords"
            buffer = [re.sub(r"^Keywords[—-]", "", line).strip()]
            continue
        heading = _ROMAN_HEADING.match(line)
        if heading:
            flush()
            section, block_type = heading.group(1).strip(), "section"
            continue
        heading = _LETTER_HEADING.match(line)
        if heading:
            flush()
            section, block_type = heading.group(1).strip(), "subsection"
            continue
        heading = _NUMBER_HEADING.match(line)
        if heading:
            flush()
            section, block_type = heading.group(1).strip(), "subsection"
            continue
        if _FIGURE.match(line):
            flush()
            old = section
            section, block_type, buffer = line.split(".", 1)[0] if "." in line else "Figure", "figure_caption", [line]
            flush("figure_caption")
            section, block_type = old, "paragraph"
            continue
        if _TABLE.match(line):
            flush()
            section, block_type, buffer = line, "table", [line]
            continue
        if _EQUATION.search(line):
            flush()
            buffer = [line]
            flush("equation")
            continue
        buffer.append(line)
        length = sum(len(item) for item in buffer)
        if scope != EvidenceScope.REFERENCE and length >= 180 and re.search(r"[.!?]$", line):
            flush("table" if block_type == "table" else "paragraph")
            block_type = "paragraph"
    flush("reference" if scope == EvidenceScope.REFERENCE else None)
    return _merge_tiny_blocks(output)


def _merge_tiny_blocks(blocks: list[PdfLogicalBlock]) -> list[PdfLogicalBlock]:
    merged: list[PdfLogicalBlock] = []
    for item in blocks:
        can_merge = item.scope == EvidenceScope.INVENTION_SOURCE and item.block_type == "paragraph" and len(item.text) < 80
        if can_merge and merged and merged[-1].page == item.page and merged[-1].scope == item.scope and merged[-1].block_type == "paragraph":
            previous = merged[-1]
            merged[-1] = PdfLogicalBlock(previous.heading, f"{previous.text} {item.text}".strip(), previous.page, previous.paragraph_index, previous.block_type, previous.scope)
        else:
            merged.append(item)
    return [PdfLogicalBlock(item.heading, item.text, item.page, index, item.block_type, item.scope) for index, item in enumerate(merged, 1)]


def read_pptx(path: Path) -> list[tuple[str, str]]:
    from pptx import Presentation
    blocks = []
    for index, slide in enumerate(Presentation(path).slides, 1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        blocks.append((f"第{index}页", text))
    return blocks
